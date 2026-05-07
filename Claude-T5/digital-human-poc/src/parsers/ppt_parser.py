"""PPT 解析器 - 提取文本、备注、图片"""

import json
from pathlib import Path
from dataclasses import dataclass, asdict

from pptx import Presentation
from pptx.util import Inches


@dataclass
class SlideContent:
    """单页 PPT 的结构化内容"""
    slide_index: int
    title: str = ""
    body_texts: list[str] = None
    notes: str = ""
    has_images: bool = False
    has_tables: bool = False
    has_charts: bool = False
    tables_data: list[list[list[str]]] = None

    def __post_init__(self):
        if self.body_texts is None:
            self.body_texts = []
        if self.tables_data is None:
            self.tables_data = []

    def to_dict(self) -> dict:
        return asdict(self)

    def to_markdown(self) -> str:
        lines = [f"## 第 {self.slide_index + 1} 页"]
        if self.title:
            lines.append(f"### {self.title}")
        for text in self.body_texts:
            lines.append(text)
        for table in self.tables_data:
            lines.append("\n| " + " | ".join(table[0]) + " |")
            lines.append("| " + " | ".join(["---"] * len(table[0])) + " |")
            for row in table[1:]:
                lines.append("| " + " | ".join(row) + " |")
        if self.notes:
            lines.append(f"\n> 备注: {self.notes}")
        return "\n".join(lines)


def extract_text_from_shape(shape) -> str:
    """从 shape 中提取纯文本"""
    if shape.has_text_frame:
        return shape.text_frame.text.strip()
    return ""


def extract_table_data(shape) -> list[list[str]]:
    """从表格 shape 中提取二维数据"""
    if not shape.has_table:
        return []
    table = shape.table
    rows = []
    for row in table.rows:
        cells = [cell.text.strip() for cell in row.cells]
        rows.append(cells)
    return rows


def parse_ppt(ppt_path: str | Path) -> list[SlideContent]:
    """解析 PPT 文件，返回每页的结构化内容"""
    ppt_path = Path(ppt_path)
    if not ppt_path.exists():
        raise FileNotFoundError(f"PPT 文件不存在: {ppt_path}")

    prs = Presentation(str(ppt_path))
    slides = []

    for i, slide in enumerate(prs.slides):
        content = SlideContent(slide_index=i)

        for shape in slide.shapes:
            # 提取标题
            try:
                pf = shape.placeholder_format
                if pf is not None and pf.idx == 1 and shape.has_text_frame:
                    content.title = shape.text_frame.text.strip()
                    continue
            except ValueError:
                pass  # 非占位符 shape，跳过

            # 提取正文文本
            text = extract_text_from_shape(shape)
            if text and text != content.title:
                content.body_texts.append(text)

            # 检测图片
            if shape.shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE
                content.has_images = True

            # 提取表格
            if shape.has_table:
                content.has_tables = True
                table_data = extract_table_data(shape)
                if table_data:
                    content.tables_data.append(table_data)

            # 检测图表
            if shape.has_chart:
                content.has_charts = True

        # 提取备注
        if slide.has_notes_slide:
            notes_text = slide.notes_slide.notes_text_frame.text.strip()
            if notes_text:
                content.notes = notes_text

        slides.append(content)

    return slides


def slides_to_markdown(slides: list[SlideContent]) -> str:
    """将所有页面内容转为 Markdown"""
    parts = [s.to_markdown() for s in slides]
    return "\n\n---\n\n".join(parts)


def slides_to_json(slides: list[SlideContent]) -> str:
    """将所有页面内容转为 JSON"""
    return json.dumps([s.to_dict() for s in slides], ensure_ascii=False, indent=2)


if __name__ == "__main__":
    import sys
    if len(sys.argv) < 2:
        print("用法: python ppt_parser.py <ppt文件路径>")
        sys.exit(1)

    result = parse_ppt(sys.argv[1])
    print(slides_to_markdown(result))

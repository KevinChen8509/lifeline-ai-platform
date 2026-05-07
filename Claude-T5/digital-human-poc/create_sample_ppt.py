"""生成一个示例 PPT 用于测试"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN


def create_sample_ppt(output_path: str = "samples/demo.pptx"):
    prs = Presentation()

    # 第1页: 标题页
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "数字人汇报系统"
    slide.placeholders[1].text = "技术方案与产品规划"

    # 第2页: 项目背景
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    slide2.shapes.title.text = "项目背景"
    body = slide2.placeholders[1].text_frame
    body.text = "传统PPT汇报需要人工演讲，效率低、一致性差"
    p2 = body.add_paragraph()
    p2.text = "AI数字人技术已趋于成熟，语音合成自然度大幅提升"
    p3 = body.add_paragraph()
    p3.text = "市场需求：企业培训、产品介绍、 investor pitch 等场景"

    # 第3页: 技术方案
    slide3 = prs.slides.add_slide(prs.slide_layouts[1])
    slide3.shapes.title.text = "核心技术方案"
    body3 = slide3.placeholders[1].text_frame
    for item in [
        "文档理解：多模态LLM解析PPT内容",
        "演讲生成：GPT-4o生成自然口语化演讲稿",
        "语音合成：Edge TTS / CosyVoice 高质量中文TTS",
        "数字人驱动：SadTalker / MuseTalk 面部动画生成",
    ]:
        p = body3.add_paragraph()
        p.text = item
        p.level = 0

    # 第4页: 产品优势
    slide4 = prs.slides.add_slide(prs.slide_layouts[1])
    slide4.shapes.title.text = "产品优势"
    body4 = slide4.placeholders[1].text_frame
    body4.text = "7×24 自动汇报，无需人工参与"
    p = body4.add_paragraph()
    p.text = "一次制作，多场景复用（汇报/培训/营销）"
    p = body4.add_paragraph()
    p.text = "支持自定义形象与音色"
    p = body4.add_paragraph()
    p.text = "成本仅为真人录制的 1/10"

    # 第5页: 总结
    slide5 = prs.slides.add_slide(prs.slide_layouts[1])
    slide5.shapes.title.text = "总结与展望"
    body5 = slide5.placeholders[1].text_frame
    body5.text = "Phase 1: PPT解析 + 演讲稿生成 + TTS合成（MVP）"
    p = body5.add_paragraph()
    p.text = "Phase 2: 多模态理解 + 情感化TTS + 口型同步"
    p = body5.add_paragraph()
    p.text = "Phase 3: 实时交互问答 + 多语言支持"
    p = body5.add_paragraph()
    p.text = "Phase 4: 3D数字人 + 企业定制化部署"

    # 添加备注
    slide2.notes_slide.notes_text_frame.text = "这里可以强调市场痛点的严重性"
    slide3.notes_slide.notes_text_frame.text = "重点讲解技术链路的完整性"

    prs.save(output_path)
    print(f"示例PPT已创建: {output_path}")


if __name__ == "__main__":
    import os
    os.makedirs("samples", exist_ok=True)
    create_sample_ppt()
